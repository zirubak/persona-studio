"""Convert a simulation transcript (markdown) into Word (.docx) + PowerPoint (.pptx).

This script is the **Python fallback** path for the docs pipeline. The preferred
orchestration (documented in `commands/simulate-*.md` Step 5) is:

    1. Try `Skill('document-skills:docx')` + `Skill('document-skills:pptx')`
       from Claude Code runtime to get Anthropic-quality outputs.
    2. If those skills are unavailable or fail, fall back to this script.

This script cannot invoke the Skill tool itself (that's a Claude Code runtime
primitive, not a Python API). So the `commands/simulate-*.md` files own
the skill-first decision; this script guarantees a working fallback.

Input contract (markdown):
- YAML frontmatter (kind, topic, participants, agenda, generated)
- H1 for document title
- H2 per agenda item + three MANDATORY close-out sections:
    - `## 결론`            conclusion (1-3 sentences, "no consensus" allowed)
    - `## 만족도 평가`      satisfaction score + per-criterion breakdown
    - `## 시스템 피드백`    persona / process / platform improvement notes
- H3 per speaker (`Lead — <name>`, `Challenger — <name>`, `결정 요약 ...`)
- blockquote lines for speech

`validate_three_axis_schema()` enforces the mandatory-sections rule. Run with
`--strict` to exit non-zero when any axis is missing. Without `--strict`, the
parser prints a warning to stderr and still produces the output files (useful
when iterating on transcript drafts).

Output directory: same as the input. Stems match the input markdown.

Usage:
    python scripts/simulation_to_docs.py <path/to/transcript.md> [--strict]

Dependencies: python-docx, python-pptx (both in pyproject dependencies).
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any


# --- frontmatter + structure parsing -------------------------------------


@dataclass
class Speaker:
    role: str           # "Lead" / "Challenger" / "결정 요약" (facilitator) / "인트로" / etc.
    name: str           # "<Persona Name>" / "퍼실리테이터"
    lines: list[str] = field(default_factory=list)   # blockquote lines (without leading ">")


@dataclass
class Section:
    title: str          # H2 text, e.g. "의제 1 — 글로벌 확장의 현재 지표와 남은 기회 구간"
    facilitator_question: str | None = None
    speakers: list[Speaker] = field(default_factory=list)


@dataclass
class Transcript:
    frontmatter: dict[str, Any]
    title: str
    intro: str               # content under the intro before first H2
    sections: list[Section] = field(default_factory=list)
    closing: Section | None = None   # "종료 요약"


_FRONTMATTER_RE = re.compile(r"^---\n(.*?)\n---\n", re.DOTALL)
_H1_RE = re.compile(r"^# (.+)$", re.MULTILINE)
_H2_RE = re.compile(r"^## (.+)$", re.MULTILINE)
_H3_RE = re.compile(r"^### (.+)$", re.MULTILINE)


def _parse_yaml_mini(text: str) -> dict[str, Any]:
    """Minimal YAML subset — key: value, key: [list], nested only one level deep.

    Avoids adding PyYAML dependency. Covers the simple frontmatter we emit.
    """
    result: dict[str, Any] = {}
    current_list_key: str | None = None
    for raw in text.splitlines():
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        if raw.startswith("  - ") and current_list_key is not None:
            # list continuation
            item = raw[4:].strip()
            item = re.sub(r"\s*#.*$", "", item).strip()
            result[current_list_key].append(item)
            continue
        current_list_key = None
        if ":" in raw:
            key, _, value = raw.partition(":")
            key = key.strip()
            value = value.strip()
            if value == "":
                result[key] = []
                current_list_key = key
            elif value.startswith("[") and value.endswith("]"):
                items = [v.strip() for v in value[1:-1].split(",") if v.strip()]
                result[key] = items
            else:
                result[key] = value
    return result


def parse_transcript(path: Path) -> Transcript:
    text = path.read_text(encoding="utf-8")

    frontmatter: dict[str, Any] = {}
    fm_match = _FRONTMATTER_RE.match(text)
    if fm_match:
        frontmatter = _parse_yaml_mini(fm_match.group(1))
        body = text[fm_match.end():]
    else:
        body = text

    h1_match = _H1_RE.search(body)
    title = h1_match.group(1).strip() if h1_match else path.stem

    # Split body by H2. Content before first H2 is intro.
    h2_positions = [(m.start(), m.end(), m.group(1).strip()) for m in _H2_RE.finditer(body)]
    if not h2_positions:
        intro = body
        return Transcript(frontmatter=frontmatter, title=title, intro=intro)

    first_h2_start = h2_positions[0][0]
    # skip H1 if present in intro area
    intro_region = body[: first_h2_start]
    if h1_match and h1_match.end() <= first_h2_start:
        intro = body[h1_match.end():first_h2_start]
    else:
        intro = intro_region
    intro = intro.strip()

    sections: list[Section] = []
    closing: Section | None = None
    for i, (_start, end, h2_text) in enumerate(h2_positions):
        next_start = h2_positions[i + 1][0] if i + 1 < len(h2_positions) else len(body)
        chunk = body[end:next_start].strip()
        section = _parse_section(h2_text, chunk)
        # treat "종료 요약" (final) as closing
        if "종료" in h2_text and "요약" in h2_text:
            closing = section
        else:
            sections.append(section)

    return Transcript(
        frontmatter=frontmatter,
        title=title,
        intro=intro,
        sections=sections,
        closing=closing,
    )


def _parse_section(h2_text: str, chunk: str) -> Section:
    section = Section(title=h2_text)

    # split by H3
    h3_positions = [(m.start(), m.end(), m.group(1).strip()) for m in _H3_RE.finditer(chunk)]

    # everything before first H3 is the facilitator question + misc prose
    if h3_positions:
        head = chunk[: h3_positions[0][0]].strip()
    else:
        head = chunk.strip()

    q_match = re.search(r"\*\*퍼실리테이터 질문[^*]*\*\*\s*:?\s*(.+?)(?=\n\n|\Z)", head, re.DOTALL)
    if q_match:
        section.facilitator_question = q_match.group(1).strip()

    for i, (_start, end, h3_text) in enumerate(h3_positions):
        next_start = h3_positions[i + 1][0] if i + 1 < len(h3_positions) else len(chunk)
        body = chunk[end:next_start].strip()
        role, name = _parse_h3(h3_text)
        speaker = Speaker(role=role, name=name, lines=_extract_blockquote_lines(body))
        section.speakers.append(speaker)
    return section


def _parse_h3(h3_text: str) -> tuple[str, str]:
    # Examples: "Lead — <Persona>", "Challenger — <Persona>", "결정 요약 (퍼실리테이터)"
    if "—" in h3_text:
        role, _, name = h3_text.partition("—")
        return role.strip(), name.strip()
    return h3_text.strip(), "퍼실리테이터"


def _extract_blockquote_lines(body: str) -> list[str]:
    """Extract continuous blockquote text. If no blockquote, use all non-empty lines."""
    lines: list[str] = []
    blockquote_seen = False
    for raw in body.splitlines():
        stripped = raw.rstrip()
        if stripped.startswith(">"):
            blockquote_seen = True
            content = stripped.lstrip(">").strip()
            lines.append(content)
        elif blockquote_seen and stripped == "":
            # blank inside blockquote: preserve paragraph boundary
            lines.append("")
    if not lines:
        lines = [line.rstrip() for line in body.splitlines() if line.strip()]
    # collapse leading/trailing empties
    while lines and not lines[0]:
        lines.pop(0)
    while lines and not lines[-1]:
        lines.pop()
    return lines


# --- Word (.docx) generation ---------------------------------------------


def build_docx(transcript: Transcript, output_path: Path) -> None:
    from docx import Document  # type: ignore
    from docx.shared import Pt  # type: ignore

    doc = Document()

    # Title
    doc.add_heading(transcript.title, level=0)

    # Metadata block
    meta = doc.add_paragraph()
    meta.add_run("회의 메타\n").bold = True
    for key in ("kind", "topic", "participants", "agenda", "generated", "facilitator"):
        if key in transcript.frontmatter:
            value = transcript.frontmatter[key]
            if isinstance(value, list):
                value = ", ".join(value)
            meta.add_run(f"- {key}: {value}\n")

    # Intro
    if transcript.intro:
        doc.add_heading("개요", level=1)
        for line in transcript.intro.split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip())

    # Sections
    for section in transcript.sections:
        doc.add_heading(section.title, level=1)
        if section.facilitator_question:
            q_para = doc.add_paragraph()
            q_run = q_para.add_run(f"[퍼실리테이터 질문] {section.facilitator_question}")
            q_run.italic = True
        for speaker in section.speakers:
            doc.add_heading(f"{speaker.role} — {speaker.name}", level=2)
            for line in speaker.lines:
                if line:
                    doc.add_paragraph(line)

    # Closing
    if transcript.closing:
        doc.add_heading("종료 요약", level=1)
        for speaker in transcript.closing.speakers:
            if speaker.role and speaker.role != "결정 요약 (퍼실리테이터)":
                doc.add_heading(speaker.role, level=2)
            for line in speaker.lines:
                if line:
                    doc.add_paragraph(line)
        # If closing was parsed as a plain section with no speakers, print raw
        if not transcript.closing.speakers:
            # Reuse the intro-style dump of the raw chunk
            pass

    # Set a consistent body font size
    for para in doc.paragraphs:
        for run in para.runs:
            if run.font.size is None:
                run.font.size = Pt(11)

    doc.save(str(output_path))


# --- PowerPoint (.pptx) generation ---------------------------------------


def build_pptx(transcript: Transcript, output_path: Path) -> None:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Title
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = transcript.title
    subtitle = slide.placeholders[1]
    sub_parts: list[str] = []
    if "topic" in transcript.frontmatter:
        sub_parts.append(f"주제: {transcript.frontmatter['topic']}")
    if "participants" in transcript.frontmatter:
        parts = transcript.frontmatter["participants"]
        if isinstance(parts, list):
            sub_parts.append(f"참석: {', '.join(parts)}")
    if "generated" in transcript.frontmatter:
        sub_parts.append(f"일시: {transcript.frontmatter['generated']}")
    subtitle.text = "\n".join(sub_parts) if sub_parts else ""

    bullet_layout = prs.slide_layouts[1]

    # Slide 2 — Agenda overview
    agenda = transcript.frontmatter.get("agenda")
    if isinstance(agenda, list) and agenda:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "의제"
        body = slide.placeholders[1].text_frame
        body.text = agenda[0]
        for item in agenda[1:]:
            body.add_paragraph().text = item

    # One slide per section — decision highlight
    for section in transcript.sections:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section.title
        body = slide.placeholders[1].text_frame
        first_line = True

        # List speakers with one-line summary (first non-empty line, truncated)
        for speaker in section.speakers:
            label = f"{speaker.role} — {speaker.name}"
            first_content = next((l for l in speaker.lines if l), "")
            summary = _truncate(first_content, 160)
            para = body if first_line else body.add_paragraph()
            if first_line:
                para.text = f"{label}: {summary}"
                first_line = False
            else:
                para.text = f"{label}: {summary}"
                para.level = 0

    # Slide — Closing consensus / actions / open questions
    if transcript.closing:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "종료 요약"
        body = slide.placeholders[1].text_frame
        first_line = True
        for line in _flatten_closing(transcript.closing):
            if first_line:
                body.text = line
                first_line = False
            else:
                body.add_paragraph().text = line

    # Slide — Next steps placeholder (parsed from action-item table if present)
    action_rows = _extract_action_items(transcript.closing)
    if action_rows:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "액션 아이템"
        body = slide.placeholders[1].text_frame
        first_line = True
        for row in action_rows:
            para = body if first_line else body.add_paragraph()
            text = f"[{row.get('담당', '')}] {row.get('내용', '')} (기한 {row.get('기한', 'TBD')})"
            if first_line:
                para.text = text
                first_line = False
            else:
                para.text = text

    prs.save(str(output_path))


def _truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def _flatten_closing(closing: Section) -> list[str]:
    out: list[str] = []
    for speaker in closing.speakers:
        for line in speaker.lines:
            if line:
                out.append(line)
    return out or [closing.facilitator_question or ""]


def _extract_action_items(closing: Section | None) -> list[dict[str, str]]:
    """Look for a Markdown table with columns 담당 / 내용 / 기한 inside the closing
    section's raw lines. Fall back to empty list if not found."""
    if closing is None:
        return []
    rows: list[dict[str, str]] = []
    for speaker in closing.speakers:
        raw_table = "\n".join(speaker.lines)
        rows.extend(_parse_md_table(raw_table))
    return rows


def _parse_md_table(text: str) -> list[dict[str, str]]:
    lines = [line.strip() for line in text.splitlines() if line.strip().startswith("|")]
    if len(lines) < 3:
        return []
    header = [c.strip() for c in lines[0].strip("|").split("|")]
    rows: list[dict[str, str]] = []
    for body_line in lines[2:]:
        cells = [c.strip() for c in body_line.strip("|").split("|")]
        if len(cells) != len(header):
            continue
        rows.append(dict(zip(header, cells)))
    return rows


# --- 3-axis schema validation --------------------------------------------


REQUIRED_AXES: tuple[str, ...] = (
    "결론",
    "만족도 평가",
    "시스템 피드백",
)


def validate_three_axis_schema(transcript: Transcript) -> list[str]:
    """Return the list of missing mandatory axes.

    Every simulation transcript must include `## 결론`, `## 만족도 평가`, and
    `## 시스템 피드백` as H2 sections (or as subsections under `## 종료 요약`).
    An empty return value means the schema is satisfied.
    """
    titles = [s.title for s in transcript.sections]
    if transcript.closing is not None:
        titles.append(transcript.closing.title)
    flat = " | ".join(titles)

    closing_body = ""
    if transcript.closing is not None:
        for speaker in transcript.closing.speakers:
            closing_body += "\n".join(speaker.lines) + "\n"

    missing: list[str] = []
    for axis in REQUIRED_AXES:
        if axis in flat:
            continue
        if axis in closing_body:
            continue
        missing.append(axis)
    return missing


# --- main -----------------------------------------------------------------


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Convert a simulation transcript (markdown) to .docx + .pptx"
    )
    parser.add_argument("input", type=Path, help="Path to simulation markdown file")
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=None,
        help="Override output directory (default: same as input file)",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Exit 2 if the 3-axis schema (결론·만족도 평가·시스템 피드백) is violated. "
        "Default: warn and continue.",
    )
    args = parser.parse_args(argv)

    if not args.input.exists():
        print(f"[ERROR] input not found: {args.input}", file=sys.stderr)
        return 2

    transcript = parse_transcript(args.input)

    missing_axes = validate_three_axis_schema(transcript)
    if missing_axes:
        msg = (
            f"[WARN] transcript missing mandatory axes: {missing_axes}. "
            "Every simulation must include `## 결론`, `## 만족도 평가`, `## 시스템 피드백`."
        )
        print(msg, file=sys.stderr)
        if args.strict:
            print("[ERROR] --strict enabled; aborting.", file=sys.stderr)
            return 2

    out_dir = args.out_dir or args.input.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = args.input.stem

    docx_path = out_dir / f"{stem}.docx"
    pptx_path = out_dir / f"{stem}.pptx"

    build_docx(transcript, docx_path)
    build_pptx(transcript, pptx_path)

    print(f"docx written: {docx_path}")
    print(f"pptx written: {pptx_path}")
    print(
        f"sections: {len(transcript.sections)}, "
        f"closing: {'yes' if transcript.closing else 'no'}, "
        f"schema: {'ok' if not missing_axes else f'missing {missing_axes}'}"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
